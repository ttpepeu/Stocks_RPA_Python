from selenium import webdriver
from selenium.webdriver.chrome.service import Service as ChromeService
from webdriver_manager.chrome import ChromeDriverManager
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import os
import shutil
import datetime

# Insert your andress of directory
directory = 'C:/your_directory_here'
NOW = datetime.datetime.now()
CSV = f'{directory}market-{NOW.month}-{NOW.year}.csv'
PP =  f'{directory}market-{NOW.month}-{NOW.year}.pptx'
CSV_BACKUP = f'{directory[2:]}market-{NOW.month}-{NOW.year}.csv'
PP_BACKUP =  f'{directory[2:]}market-{NOW.month}-{NOW.year}.pptx'
now = datetime.datetime.now()
DRIVER = webdriver.Chrome(service=ChromeService(ChromeDriverManager().install()))
d = DRIVER

class Main:
    def Backup():
        if (os.path.exists(CSV_BACKUP)):
            shutil.move(CSV_BACKUP,f'{directory[2:]}Backup')
            shutil.move(PP_BACKUP,f'{directory[2:]}Backup')

    def running():
        class GetDatas:
            def __init__(self,links):
                self.links = links

            def webscraping(self):
                d.get(self.links)
                company = d.find_element('xpath','//*[@id="rcnt"]/div[2]/div/div/div[3]/div[1]/div/div/div[2]/div[2]/div[1]/div/span').text[6:]
                value = d.find_element('xpath','//*[@id="knowledge-finance-wholepage__entity-summary"]/div[3]/g-card-section/div/g-card-section/div[2]/div[1]/span[1]/span/span[1]').text
                coin = d.find_element('xpath','//*[@id="knowledge-finance-wholepage__entity-summary"]/div[3]/g-card-section/div/g-card-section/div[2]/div[1]/span[1]/span/span[2]').text
                
                return company, value, coin
            
            
        companies = []
        values = []
        coins = []

        # PETR4, Alphabet BDR, APPL34, AMZO34
        array_companies = [
            'https://www.google.com/search?client=opera&q=ações+petrobras&sourceid=opera&ie=UTF-8&oe=UTF-8',
            'https://www.google.com/search?q=ações+google&client=opera&hs=tuO&sxsrf=APwXEdf6fBcfT8nu3VIh24jLqaltXjGucw%3A1685455752875&ei=iAN2ZPuKNfOh1sQPusezyAY&ved=0ahUKEwj77rz7m53_AhXzkJUCHbrjDGkQ4dUDCA8&uact=5&oq=ações+google&gs_lcp=Cgxnd3Mtd2l6LXNlcnAQAzIQCAAQgAQQsQMQgwEQRhD6ATIFCAAQgAQyBQgAEIAEMgUIABCABDIFCAAQgAQyBQgAEIAEMgUIABCABDIFCAAQgAQyBQgAEIAEMgUIABCABDoKCAAQRxDWBBCwAzoKCAAQigUQsAMQQzoLCAAQgAQQsQMQgwE6CAgAEIAEELEDSgQIQRgAUMW0nQFY6LmdAWCLu50BaANwAXgAgAF4iAGWBZIBAzEuNZgBAKABAcABAcgBCg&sclient=gws-wiz-serp',
            'https://www.google.com/search?q=ações+apple&client=opera&sxsrf=APwXEdc4otC_a_tTKeNPUBqu6kmR0uHzXg%3A1685458904321&ei=2A92ZMKgE7PZ5OUP4Le5qA0&ved=0ahUKEwjC7Jnap53_AhWzLLkGHeBbDtUQ4dUDCA8&uact=5&oq=ações+apple&gs_lcp=Cgxnd3Mtd2l6LXNlcnAQAzIQCAAQgAQQsQMQgwEQRhD6ATIFCAAQgAQyBQgAEIAEMgUIABCABDIFCAAQgAQyBQgAEIAEMgUIABCABDIFCAAQgAQyBQgAEIAEMgUIABCABDoKCAAQRxDWBBCwAzoKCCMQigUQJxCdAjoPCCMQigUQJxCdAhBGEPoBOgsIABCABBCxAxCDAToICAAQgAQQsQM6BwgAEIoFEEM6BwgjEIoFECc6DQgAEIoFELEDEIMBEENKBAhBGABQmgRY2RFghhRoAXAAeAGAAcsCiAHjDJIBCDAuMTEuMC4xmAEAoAEBwAEByAEI&sclient=gws-wiz-serp',
            'https://www.google.com/search?q=ações+amazon&client=opera&sxsrf=APwXEdfCQviIojRzQQe0zpk3oV14tckXkg%3A1685459159893&ei=1xB2ZKabNqyz5OUPnIKPuAQ&oq=ações+amaxo&gs_lcp=Cgxnd3Mtd2l6LXNlcnAQAxgAMhIIABANEIAEELEDEIMBEEYQ-gEyBwgAEA0QgAQyBwgAEA0QgAQyBwgAEA0QgAQyBwgAEA0QgAQyBwgAEA0QgAQyBwgAEA0QgAQyBwgAEA0QgAQyBwgAEA0QgAQyBwgAEA0QgAQ6CggAEEcQ1gQQsAM6CggAEIoFELADEEM6CggjEIoFECcQnQI6BQgAEIAEOg8IIxCKBRAnEJ0CEEYQ-gE6CwgAEIAEELEDEIMBOg0IABCKBRCxAxCDARBDOggIABCABBCxAzoHCCMQigUQJzoHCAAQigUQQzoNCAAQgAQQsQMQgwEQCjoKCAAQgAQQsQMQCjoLCAAQFhAeEPEEEApKBAhBGABQvQZYxRNg9ypoAnABeAGAAecCiAGGCpIBBzAuOC4wLjGYAQCgAQHAAQHIAQo&sclient=gws-wiz-serp'
        ]


        for i in array_companies:
            data = GetDatas(i).webscraping()
            company, value, coin = data
            companies.append(company)
            values.append(value)
            coins.append(coin)

        df = pd.DataFrame({'company':companies, 'value':values, 'type_coin':coins})

        df.to_csv(CSV, index=False)

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        rows = len(array_companies)+1
        cols = len(df.iloc[0])
        x = Inches(2.5)
        y = Inches(2.0)
        width = Inches(5.0)
        height = Inches(1.0)
        table = slide.shapes.add_table(rows,cols,x,y,width,height).table
        title.text = "Summary of Market"

        for i in range(cols):
            table.cell(0,i).text=str(df.columns[i])

        for i in range(len(array_companies)):
            for v in range(cols):
                table.cell(i+1,v).text=str(df.iloc[i,v])

        prs.save(PP)
        d.quit()
        
if __name__ == '__main__':
    while True:
        try:
            Main.Backup()
        except:
            print(f'Already there the file: {CSV}')
            print(f'Already there the file: {PP}')
            break
        Main.running()
    d.close()

